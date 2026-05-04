import os
import sys
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PPTXGeneratorV2:
    def __init__(self, output_path):
        self.prs = Presentation()
        # Set slide size to 16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.output_path = output_path
        
        # Colors (Premium Navy Theme)
        self.NAVY = RGBColor(0x15, 0x4B, 0x98)
        self.WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        self.CYAN = RGBColor(0x22, 0xD3, 0xEE)
        self.GRAY_BG = RGBColor(0xF3, 0xF4, 0xF6)
        self.BLACK = RGBColor(0x00, 0x00, 0x00)

    def set_font(self, run, size, bold=False, color=None):
        run.font.name = 'Arial'
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    def add_title_slide(self, title, subtitle):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.NAVY
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self.set_font(p.runs[0], 44, bold=True, color=self.WHITE)
        
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
        tf2 = txBox2.text_frame
        tf2.text = subtitle
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        self.set_font(p2.runs[0], 24, color=self.CYAN)

    def add_content_slide(self, title, content_list, bg_color=None, is_prompt=False):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5]) # Title Only
        
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self.set_font(title_shape.text_frame.paragraphs[0].runs[0], 32, bold=True, color=self.NAVY)
        
        # Content box
        top = Inches(1.5)
        left = Inches(0.5)
        width = Inches(12.33)
        height = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            if is_prompt:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.JUSTIFY
            
            p.space_after = Pt(8)
            
            if item.startswith("-"):
                p.level = 1
                text = item[1:].strip()
                p.text = text
                self.set_font(p.runs[0], 18, color=self.BLACK)
            elif item.startswith("###") or item.startswith("##"):
                p.text = item.replace("#", "").strip()
                self.set_font(p.runs[0], 22, bold=True, color=self.NAVY)
            else:
                self.set_font(p.runs[0], 19, color=self.BLACK)

    def save(self):
        self.prs.save(self.output_path)
        print(f"Presentation saved to {self.output_path}")

def parse_markdown_to_slides(file_path):
    """Simple parser for pedagogical markdown drafts."""
    slides = []
    current_slide = None
    
    if not os.path.exists(file_path):
        return []

    with open(file_path, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line: continue
            
            if line.startswith("# "):
                title = line.replace("#", "").strip()
                slides.append({"type": "title", "title": title, "subtitle": "NoteBookLLM Pedagogical Export"})
            elif line.startswith("## "):
                if current_slide: slides.append(current_slide)
                current_slide = {"type": "content", "title": line.replace("##", "").strip(), "content": []}
            elif current_slide:
                current_slide["content"].append(line)
        
        if current_slide: slides.append(current_slide)
    return slides

def main():
        "1. 'I've been living here for five years.' (Nhấn mạnh sự liên tục).",
        "2. 'It’s been five years since I moved here.' (Cấu trúc thông dụng)."
    ])

    # CONCLUSION
    gen.add_content_slide("Công thức dẫn đầu: ROLE + TASK + CONTEXT", [
        "Để học tập hiệu quả cùng AI, hãy đổi từ 'tra cứu' sang 'đối thoại'.",
        "1. ROLE: Giao vai trò chuyên gia cho AI.",
        "2. TASK: Mô tả nhiệm vụ cụ thể, rõ ràng.",
        "3. CONTEXT: Cung cấp bối cảnh (cấp học, mục tiêu).",
        "## Hãy bắt đầu RTC ngay hôm nay!"
    ])

    gen.save()

if __name__ == "__main__":
    main()
