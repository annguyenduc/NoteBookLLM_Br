import os
import json
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class HTMLToPPTXSync:
    def __init__(self, output_path):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.output_path = output_path
        
        # Design System Tokens (Sync with HTML CSS)
        self.NAVY = RGBColor(0x00, 0x1F, 0x3F)
        self.WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        self.BLACK = RGBColor(0x33, 0x33, 0x33)
        self.ROLE_BLUE = RGBColor(0x34, 0x98, 0xDB)
        self.TASK_GREEN = RGBColor(0x2E, 0xCC, 0x71)
        self.CONTEXT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
        self.ERROR_RED = RGBColor(0xE7, 0x4C, 0x3C)
        self.GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)

    def _setup_blank_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.NAVY
        bg.line.width = 0
        return slide

    def _add_shadow(self, shape):
        # Basic shadow since complex color/alpha is not available via simple API
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.blur_radius = Pt(12)
        shadow.distance = Pt(5)
        shadow.angle = 45

    def _add_text_to_frame(self, frame, text, font_size, italic=False, bold=False, color=None, alignment=PP_ALIGN.LEFT):
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = frame.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.italic = italic
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return run

    def add_title(self, slide, soup):
        title_tag = soup.find('h1')
        title_text = title_tag.get_text() if title_tag else ""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
        self._add_text_to_frame(title_box.text_frame, title_text.upper(), Pt(44), bold=True, color=self.WHITE, alignment=PP_ALIGN.CENTER)
        title_box.text_frame.paragraphs[0].runs[0].font.name = 'Montserrat'

        subtitle_tag = soup.find('div', class_='subtitle')
        if subtitle_tag:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.33), Inches(0.4))
            self._add_text_to_frame(sub_box.text_frame, subtitle_tag.get_text(), Pt(18), color=self.WHITE, alignment=PP_ALIGN.CENTER)

    def _add_rich_text(self, frame, html_box, base_font_size, base_color=None, italic=False):
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.JUSTIFY
        
        # Mapping colors from classes
        color_map = {
            'highlight-r': self.ROLE_BLUE,
            'highlight-t': self.TASK_GREEN,
            'highlight-c': self.CONTEXT_ORANGE
        }

        # Iterating through contents for basic rich text
        for child in html_box.descendants:
            if isinstance(child, str):
                if child.parent.name in ['strong', 'b']:
                    run = p.add_run()
                    run.text = child
                    run.font.bold = True
                elif child.parent.name == 'span':
                    run = p.add_run()
                    run.text = child
                    cls = child.parent.get('class', [None])[0]
                    if cls in color_map:
                        run.font.color.rgb = color_map[cls]
                        run.font.bold = True
                else:
                    # Normal text - check if it's already inside a strong/span handled above
                    if child.parent.name not in ['strong', 'b', 'span']:
                        run = p.add_run()
                        run.text = child
            
            # Apply base formatting to all runs in this paragraph
            for run in p.runs:
                run.font.size = base_font_size
                if italic: run.font.italic = True
                # Set base color if no specific color was applied (Role/Task/Context etc)
                try:
                    _ = run.font.color.rgb
                except (AttributeError, ValueError):
                    if base_color:
                        run.font.color.rgb = base_color

    def sync_anatomy_slide(self, html_path):
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')

        slide = self._setup_blank_slide()
        self.add_title(slide, soup)

        # Main Container
        glass = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.6))
        glass.fill.solid()
        glass.fill.fore_color.rgb = self.WHITE
        glass.line.width = 0
        self._add_shadow(glass)

        # Parse Blocks
        blocks = soup.find_all('div', class_='lego-block')
        colors = [self.ROLE_BLUE, self.TASK_GREEN, self.CONTEXT_ORANGE]
        
        base_y = 1.9
        for i, block in enumerate(blocks):
            label = block.find('span', class_='label').get_text()
            content_tag = block.find('div', class_='content')
            header = content_tag.find('strong').get_text().strip()
            detail = content_tag.find('small').get_text().strip()
            color = colors[i]

            # LEGO Nub
            nub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), Inches(base_y - 0.1), Inches(0.3), Inches(0.2))
            nub.fill.solid()
            nub.fill.fore_color.rgb = color
            nub.line.width = 0

            # Block Shape
            blk = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(base_y), Inches(10.9), Inches(1.0))
            blk.fill.solid()
            blk.fill.fore_color.rgb = color
            blk.line.width = 0

            # Label Badge
            bgc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.4), Inches(base_y + 0.15), Inches(0.7), Inches(0.7))
            bgc.fill.solid()
            bgc.fill.fore_color.rgb = self.WHITE
            bgc.text = label
            brun = bgc.text_frame.paragraphs[0].runs[0]
            brun.font.bold = True
            brun.font.size = Pt(28)
            brun.font.color.rgb = color

            # Text
            txt = slide.shapes.add_textbox(Inches(2.4), Inches(base_y + 0.1), Inches(9.5), Inches(0.8))
            tf = txt.text_frame
            self._add_text_to_frame(tf, header, Pt(20), bold=True, color=self.WHITE)
            
            p2 = tf.add_paragraph()
            p2.text = detail
            p2.runs[0].font.size = Pt(14)
            p2.runs[0].font.color.rgb = self.WHITE

            base_y += 1.3

        # Full Prompt Box
        prompt_box = soup.find('div', class_='prompt-box')
        if prompt_box:
            pb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.8), Inches(10.9), Inches(1.0))
            pb.fill.solid()
            pb.fill.fore_color.rgb = self.GRAY_BG
            pb.line.color.rgb = self.NAVY
            pb.line.width = Pt(1.5)
            
            pbt = slide.shapes.add_textbox(Inches(1.4), Inches(5.9), Inches(10.5), Inches(0.8))
            self._add_rich_text(pbt.text_frame, prompt_box, Pt(14), italic=True, base_color=self.BLACK)

    def sync_battle_slide(self, html_path):
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')

        slide = self._setup_blank_slide()
        self.add_title(slide, soup)

        cards = soup.find_all('div', class_='card')
        x_offsets = [0.5, 6.8]

        for i, card in enumerate(cards):
            cx = Inches(x_offsets[i])
            is_good = 'good' in card['class']
            color = self.TASK_GREEN if is_good else self.ERROR_RED
            header_text = card.find('div', class_='card-header').get_text(strip=True)
            badge_text = card.find('span', class_='badge').get_text(strip=True)
            prompt_text = card.find('div', class_='prompt-text').get_text(strip=True)
            result_text = card.find('div', class_='result-box').get_text(strip=True)
            
            # Card logic: Use two overlapping rounded rectangles to simulate border-top-radius
            # 1. The "Accent" (Background layer with curve)
            accent_h = Inches(1.0) # Only need the top part
            acc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.5), Inches(6.0), accent_h)
            acc.fill.solid()
            acc.fill.fore_color.rgb = color
            acc.line.width = 0
            # Set corner radius (0.2 is common for 16px equivalent)
            acc.adjustments[0] = 0.15 
            
            # 2. The "Body" (White layer on top, shifted down slightly)
            body_y_offset = Pt(6) # Thickness of the border-top
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.5) + body_y_offset, Inches(6.0), Inches(5.2) - body_y_offset)
            sh.fill.solid()
            sh.fill.fore_color.rgb = self.WHITE
            sh.line.width = 0
            sh.adjustments[0] = 0.15
            self._add_shadow(sh)

            # Badge & Header
            txt = slide.shapes.add_textbox(cx + Inches(0.2), Inches(1.7), Inches(5.6), Inches(0.5))
            tf = txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = f" {badge_text} "
            r1.font.bold = True
            r1.font.size = Pt(14)
            r1.font.color.rgb = color
            
            r2 = p.add_run()
            r2.text = f"   {header_text.replace(badge_text, '').strip()}"
            r2.font.bold = True
            r2.font.size = Pt(18) # Giảm 1 chút để tránh tràn header
            r2.font.color.rgb = self.NAVY

            # Prompt Area
            pbox_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.2), Inches(2.4), Inches(5.6), Inches(1.6))
            pbox_shape.fill.solid()
            pbox_shape.fill.fore_color.rgb = self.GRAY_BG
            pbox_shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            pbox_shape.line.width = Pt(1)
            
            pbt = slide.shapes.add_textbox(cx + Inches(0.3), Inches(2.5), Inches(5.4), Inches(1.4))
            # Clean double quotes from HTML if any, then wrap in single ones
            clean_prompt = card.find('div', class_='prompt-text').get_text(strip=True).strip('"')
            self._add_rich_text(pbt.text_frame, card.find('div', class_='prompt-text'), Pt(13), italic=True, base_color=self.BLACK)

            # Result Area
            rbox_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.2), Inches(4.2), Inches(5.6), Inches(1.0))
            rbox_shape.fill.solid()
            rbox_shape.fill.fore_color.rgb = RGBColor(0xFE, 0xE2, 0xE2) if not is_good else RGBColor(0xDC, 0xFC, 0xE7)
            rbox_shape.line.width = 0
            
            rt = slide.shapes.add_textbox(cx + Inches(0.3), Inches(4.25), Inches(5.4), Inches(0.9))
            self._add_rich_text(rt.text_frame, card.find('div', class_='result-box'), Pt(11), base_color=self.BLACK)

            # Stats (Scraping Labels)
            stat_bars = card.find_all('div', class_='stat-bar')
            for j, sbar in enumerate(stat_bars):
                # HTML structure: [label_span, bar_div, value_span] - but let's just get all text
                spans = sbar.find_all('span')
                label = spans[0].get_text() if spans else ""
                value = spans[1].get_text() if len(spans) > 1 else ""
                
                percent = 1.0 if (is_good and j==1) or (not is_good and j==0) else 0.1
                if is_good and j==0: percent = 0.9 # Length ~40/45
                
                sy = 5.4 + (j * 0.5)
                self._add_stat_v5(slide, cx + Inches(0.2), Inches(sy), label, value, percent, color)

    def _add_stat_v5(self, slide, x, y, label, value, percent, color):
        # Label Left
        lbl = slide.shapes.add_textbox(x, y, Inches(1.4), Inches(0.3))
        self._add_text_to_frame(lbl.text_frame, label, Pt(10), bold=True, color=self.NAVY)

        # Bar BG
        bar_x = x + Inches(1.5)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, y + Inches(0.08), Inches(3.0), Inches(0.12))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        bg.line.width = 0
        
        # Bar Fill
        fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, y + Inches(0.08), Inches(3.0 * percent), Inches(0.12))
        fill.fill.solid()
        fill.fill.fore_color.rgb = color
        fill.line.width = 0

        # Value Right
        vlbl = slide.shapes.add_textbox(bar_x + Inches(3.1), y, Inches(1.2), Inches(0.3))
        self._add_text_to_frame(vlbl.text_frame, value, Pt(10), bold=True, color=self.NAVY)

    def save(self):
        self.prs.save(self.output_path)
        print(f"File PPTX V5 (Final Polish) saved: {self.output_path}")

if __name__ == "__main__":
    output_file = "d:\\NoteBookLLM_Br\\assets\\Math10_Prompt_Engineering_Premium_V4.pptx"
    sync = HTMLToPPTXSync(output_file)
    sync.sync_anatomy_slide("d:\\NoteBookLLM_Br\\assets\\RTC_Anatomy_Slide.html")
    sync.sync_battle_slide("d:\\NoteBookLLM_Br\\assets\\RTC_Length_Battle_Slide.html")
    sync.save()
