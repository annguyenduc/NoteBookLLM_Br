import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PPTXGeneratorV2:
    def __init__(self, output_path):
        self.prs = Presentation()
        # Set slide size to 16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.output_path = output_path
        
        # Colors
        self.NAVY = RGBColor(0x15, 0x4B, 0x98)
        self.WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        self.CYAN = RGBColor(0x22, 0xD3, 0xEE)
        self.GRAY_BG = RGBColor(0xF3, 0xF4, 0xF6)
        self.BLACK = RGBColor(0x00, 0x00, 0x00)

    def set_font(self, run, size, bold=False, color=None):
        run.font.name = 'Arial'
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    def add_title_slide(self, title, subtitle):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.NAVY
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self.set_font(p.runs[0], 44, bold=True, color=self.WHITE)
        
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
        tf2 = txBox2.text_frame
        tf2.text = subtitle
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        self.set_font(p2.runs[0], 24, color=self.CYAN)

    def add_content_slide(self, title, content_list, bg_color=None, is_prompt=False):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5]) # Title Only
        
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self.set_font(title_shape.text_frame.paragraphs[0].runs[0], 32, bold=True, color=self.NAVY)
        
        # Content box
        top = Inches(1.5)
        left = Inches(0.5)
        width = Inches(12.33)
        height = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            if is_prompt:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.JUSTIFY
            
            p.space_after = Pt(8)
            
            if item.startswith("-"):
                p.level = 1
                text = item[1:].strip()
                p.text = text
                self.set_font(p.runs[0], 18, color=self.BLACK)
            elif item.startswith("###") or item.startswith("##"):
                p.text = item.replace("#", "").strip()
                self.set_font(p.runs[0], 22, bold=True, color=self.NAVY)
            else:
                self.set_font(p.runs[0], 19, color=self.BLACK)

    def save(self):
        self.prs.save(self.output_path)
        print(f"V2 Presentation saved to {self.output_path}")

def main():
    output = "d:\\NoteBookLLM_Br\\brain\\distilled\\ATOMS_Demo_Prompt_Comparison_K10_V2.pptx"
    gen = PPTXGeneratorV2(output)
    
    # Slides mapping according to presentation_draft.md (Full Content version)
    gen.add_title_slide("KỸ THUẬT PROMPT RTC K10", "Dành cho việc Giảng dạy & Đối soát Đa môn (Full Content)")
    
    # Theory
    gen.add_content_slide("Tại sao cần RTC Management?", [
        "AI sở hữu kho tàng tri thức khổng lồ nhưng hỗn tạp.",
        "Prompt không cấu trúc (Simple) chỉ kích hoạt tầng mặt của tri thức.",
        "Prompt cấu trúc (RTC) giúp AI 'đóng vai' chuyên gia, kích hoạt đúng vùng sư phạm chuyên sâu.",
        "## Cấu thức: ROLE + TASK + CONTEXT"
    ])
    
    # TOAN
    gen.add_content_slide("[TOÁN] Lượt A - Simple Prompt", [
        "## Yêu cầu (Prompt):",
        "Giải phương trình x² - 5x + 6 = 0 và tính tổng, tích các nghiệm.",
        "## Phản hồi AI:",
        "Cách 1: Phân tích thành nhân tử (x - 2)(x - 3) = 0 => x1 = 2, x2 = 3.",
        "Cách 2: Delta (Δ) = 1 => x1 = 3, x2 = 2.",
        "S = 5, P = 6."
    ], bg_color=gen.GRAY_BG)
    
    gen.add_content_slide("[TOÁN] Lượt B - RTC (Phần 1/2)", [
        "## Prompt cấu trúc:",
        "ROLE: Giáo viên Toán K-12... TASK: Giải & chứng minh bằng Vi-ét...",
        "## Phản hồi Sư phạm từ AI:",
        "Chào các em, hôm nay chúng ta xét phương trình: x² - 5x + 6 = 0",
        "### 1. Giải phương trình bằng Delta (Δ)",
        "Bước 1: Hệ số a = 1, b = -5, c = 6",
        "Bước 2: Tính Δ = b² - 4ac = (-5)² - 4*1*6 = 25 - 24 = 1",
        "Bước 3: Tìm nghiệm x1, x2: x1 = 3, x2 = 2. Vậy S = {2; 3}."
    ])
    
    gen.add_content_slide("[TOÁN] Lượt B - RTC (Phần 2/2)", [
        "### 2. Áp dụng hệ thức Vi-ét kiểm tra",
        "- Tổng S = x1 + x2 = 3 + 2 = 5. (Theo lý thuyết -b/a = 5) -> Khớp!",
        "- Tích P = x1 * x2 = 3 * 2 = 6. (Theo lý thuyết c/a = 6) -> Khớp!",
        "### 3. Kết luận",
        "Nghiệm của phương trình là x = 2 và x = 3 hoàn toàn chính xác."
    ])

    # VAT LY
    gen.add_content_slide("[VẬT LÝ] Lượt A - Simple Prompt", [
        "## Prompt:",
        "Giải bài tập xe chuyển động thẳng nhanh dần đều (v0=10, a=2).",
        "## Response:",
        "v = 10 + 2t; s = 10t + t².",
        "Sau 5 giây: s = 75m."
    ], bg_color=gen.GRAY_BG)
    
    gen.add_content_slide("[VẬT LÝ] Lượt B - RTC (Phần 1/3)", [
        "## Phân tích Sư phạm:",
        "Chào em, thầy rất vui được đồng hành cùng em. Đây là một bài toán cơ bản giúp em nắm vững động học.",
        "### (1) Tóm tắt đề bài",
        "- Vận tốc đầu (v0): 10 m/s",
        "- Gia tốc (a): 2 m/s²",
        "- Thời gian (t): 5 s",
        "- Yêu cầu: Tính quãng đường (s)."
    ])
    
    gen.add_content_slide("[VẬT LÝ] Lượt B - RTC (Phần 2/3)", [
        "### (2) Phân tích hiện tượng",
        "- Loại chuyển động: Vì a=2 (>0) và v0=10 (>0), nên đây là chuyển động THẲNG NHANH DẦN ĐỀU.",
        "- Diễn biến: Cứ mỗi giây, vận tốc tăng thêm 2m/s. Quãng đường các giây sau dài hơn giây trước.",
        "- Chiều dương: Chọn chiều chuyển động của xe."
    ])
    
    gen.add_content_slide("[VẬT LÝ] Lượt B - RTC (Phần 3/3)", [
        "### (3) Các bước giải bài toán",
        "Bước 1: Chọn hệ quy chiếu gốc thời gian lúc xe bắt đầu.",
        "Bước 2: Áp dụng công thức: s = v0.t + 0.5.a.t²",
        "Bước 3: Thay số: s = 10*5 + 0.5*2*25 = 75 (m).",
        "## Kết luận: Sau 5 giây, xe đi được 75 mét."
    ])

    # HOA HOC
    gen.add_content_slide("[HÓA HỌC] Lượt A - Simple Prompt", [
        "Cấu hình Magnesium (Z=12): 1s² 2s² 2p⁶ 3s².",
        "Vị trí: Ô 12, Chu kỳ 3, Nhóm IIA."
    ], bg_color=gen.GRAY_BG)
    
    gen.add_content_slide("[HÓA HỌC] Lượt B - RTC (Phần 1/2)", [
        "### 1. Cấu tạo nguyên tử Magnesium (Z = 12)",
        "Chào em, cùng thầy tìm hiểu về nguyên tử này nhé:",
        "- Hạt nhân: 12 Proton. Vỏ: 12 Electron.",
        "- Phân bố Electron:",
        "  - Lớp 1 (trong cùng): 2 e.",
        "  - Lớp 2: 8 e (tối đa).",
        "  - Lớp 3 (ngoài cùng): 2 e.",
        "- Cấu hình đầy đủ: 1s² 2s² 2p⁶ 3s²."
    ])
    
    gen.add_content_slide("[HÓA HỌC] Lượt B - RTC (Phần 2/2)", [
        "### 2. Xác định vị trí trong Bảng tuần hoàn",
        "- Ô nguyên tố: 12 (do Z = 12).",
        "- Chu kỳ: 3 (do có 3 lớp electron).",
        "- Nhóm: IIA (do có 2 e lớp ngoài cùng - nguyên tố s).",
        "## Tóm tắt:",
        "Magnesium là kim loại mạnh, thuộc nhóm kim loại kiềm thổ. Mẹo: Tìm hàng 3, cột 2 trên bảng hệ thống."
    ])

    # NGU VAN
    gen.add_content_slide("[NGỮ VĂN] Lượt A - Simple Prompt", [
        "Dàn ý mạng xã hội: Mở bài (giới thiệu), Thân bài (lợi ích/tác hại), Kết bài (liên hệ)."
    ], bg_color=gen.GRAY_BG)
    
    gen.add_content_slide("[NGV] RTC: Tư duy Phản biện (Part 1/3)", [
        "### I. Mở bài: Dẫn dắt vấn đề",
        "- Hook: 'Cảnh tượng những nhóm học sinh mỗi người dán mắt vào một màn hình không còn xa lạ...'",
        "- Luận điểm: Mạng xã hội là 'con dao hai lưỡi' tì thuộc vào bản lĩnh người dùng."
    ])
    
    gen.add_content_slide("[NGV] RTC: Đối lập đa chiều (Part 2/3)", [
        "### II. Thân bài",
        "- 1. Giải thích: Lan tỏa nhanh, không biên giới, đầy tin chưa kiểm chứng.",
        "- 2. Tác động tích cực: Kho tàng tri thức, kết nối đam mê, rèn kỹ năng số.",
        "- 3. Tác động tiêu cực: Xao nhãng học tập, FOMO, Áp lực đồng lứa.",
        "- 4. Nguyên nhân: Thuật toán gây nghiện & thiếu kỹ năng quản lý thời gian."
    ])
    
    gen.add_content_slide("[NGV] RTC: Thông điệp (Part 3/3)", [
        "### 5. Bác bỏ & Mở rộng",
        "- Phản biện: Đừng cấm đoán, hãy giáo dục về 'Năng lực số' (Digital Literacy).",
        "### III. Kết bài",
        "- Thông điệp: Hãy là người dùng mạng xã hội 'tỉnh thức'.",
        "- 'Hãy tắt màn hình để bật lên những kết nối thật...'"
    ])

    # TIENG ANH
    gen.add_content_slide("[TIẾNG ANH] Lượt A - Simple Prompt", [
        "Lỗi: 'I have live here since 5 years.'",
        "Sửa: 'I have lived here for 5 years.'"
    ], bg_color=gen.GRAY_BG)
    
    gen.add_content_slide("[ANH] RTC: Chuyên gia ngôn ngữ (Part 1/2)", [
        "### 1. Câu sửa lỗi: 'I have lived here for 5 years.'",
        "### 2. Giải thích lỗi sai",
        "- Lỗi 1 (Chia động từ): Sai 'have live' -> Đúng 'have lived' (Have/Has + V3).",
        "- Lỗi 2 (Giới từ): Sai 'since 5 years' -> Đúng 'for 5 years' (Dùng cho khoảng thời gian)."
    ])
    
    gen.add_content_slide("[ANH] RTC: Chuyên gia ngôn ngữ (Part 2/2)", [
        "### 3. Phân biệt Past Simple vs. Present Perfect",
        "- Present Perfect: Bắt đầu quá khứ, hiện tại vẫn tiếp tục.",
        "### 4. Gợi ý Natural Alternatives (Cách nói tự nhiên)",
        "1. 'I've been living here for five years.' (Nhấn mạnh sự liên tục).",
        "2. 'It’s been five years since I moved here.' (Cấu trúc thông dụng)."
    ])

    # CONCLUSION
    gen.add_content_slide("Công thức dẫn đầu: ROLE + TASK + CONTEXT", [
        "Để học tập hiệu quả cùng AI, hãy đổi từ 'tra cứu' sang 'đối thoại'.",
        "1. ROLE: Giao vai trò chuyên gia cho AI.",
        "2. TASK: Mô tả nhiệm vụ cụ thể, rõ ràng.",
        "3. CONTEXT: Cung cấp bối cảnh (cấp học, mục tiêu).",
        "## Hãy bắt đầu RTC ngay hôm nay!"
    ])

    gen.save()

if __name__ == "__main__":
    main()
